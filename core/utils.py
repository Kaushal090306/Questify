import fitz  # PyMuPDF
import docx
import openai
import json
import os
from pptx import Presentation
import re
from django.conf import settings

openai.api_key = os.getenv("OPENAI_API_KEY", "")

def parse_document(file):
    """
    Extracts and returns text from the uploaded file (PDF, DOCX, PPTX).
    """
    file_ext = file.name.lower().split('.')[-1]
    raw_text = ""

    if file_ext == "pdf":
        # Parse PDF
        pdf_doc = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf_doc:
            raw_text += page.get_text()
        pdf_doc.close()
    elif file_ext == "docx":
        # Parse DOCX
        doc = docx.Document(file)
        for para in doc.paragraphs:
            raw_text += para.text + "\n"
    elif file_ext == "pptx":
        # Parse PPTX
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text += shape.text + "\n"
                    
            # Also extract notes from slides
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, "text"):
                        raw_text += shape.text + "\n"
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, DOCX, or PPTX.")

    return raw_text.strip()

def generate_quiz_from_text(text, num_questions=5, difficulty='medium', quiz_type='multiple_choice'):
    """
    Generates a quiz from the given text using OpenAI. 
    Supports multiple question types including descriptive and fill-in-the-blank.
    Returns a list of question dicts as specified.
    """
    try:
        # Clean the text
        cleaned_text = clean_text(text)
        
        # Limit text length to avoid token limits
        max_chars = 8000
        if len(cleaned_text) > max_chars:
            cleaned_text = cleaned_text[:max_chars] + "..."
        
        # Create the prompt based on question type
        prompt = create_enhanced_quiz_prompt(cleaned_text, num_questions, difficulty, quiz_type)
        
        # Call OpenAI API
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert quiz generator. Generate educational quiz questions based on the provided text. Always return valid JSON format."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=3000
        )
        
        # Parse the response
        content = response.choices[0].message.content
        
        # Extract JSON from response
        quiz_data = extract_json_from_response(content)
        
        # Validate and clean the quiz data
        validated_quiz = validate_enhanced_quiz_data(quiz_data, num_questions, quiz_type)
        
        return validated_quiz
        
    except Exception as e:
        raise ValueError(f"Error generating quiz: {str(e)}")

def create_enhanced_quiz_prompt(text, num_questions, difficulty, quiz_type):
    """Create enhanced prompt for multiple question types"""
    
    difficulty_instructions = {
        'easy': 'Create simple, straightforward questions that test basic understanding.',
        'medium': 'Create moderately challenging questions that require comprehension and analysis.',
        'hard': 'Create complex questions that require deep understanding and critical thinking.'
    }
    
    type_instructions = {
        'multiple_choice': 'Each question should have 4 options (A, B, C, D) with exactly one correct answer.',
        'true_false': 'Each question should be a statement that can be answered with True or False.',
        'fill_in_blank': 'Each question should have a missing word or phrase that needs to be filled in. Use ___ to indicate the blank.',
        'descriptive': 'Each question should require a detailed written response (2-3 sentences).',
        'mixed': 'Create a mix of multiple choice, fill-in-the-blank, and descriptive questions.'
    }
    
    base_prompt = f"""
Based on the following text, generate {num_questions} {difficulty} difficulty quiz questions.

DIFFICULTY LEVEL: {difficulty_instructions[difficulty]}

QUESTION TYPE: {type_instructions[quiz_type]}

REQUIREMENTS:
1. Questions must be directly answerable from the provided text
2. Ensure questions cover different parts of the text
3. For multiple choice: Make distractors plausible but clearly incorrect
4. For fill-in-the-blank: Choose important keywords or phrases to remove
5. For descriptive: Ask questions that require explanation or analysis
6. Include explanations for all answers
7. Assign appropriate topics/categories to each question

Return your response in this exact JSON format:
{{
  "quiz": [
    {{
      "question_text": "Your question here",
      "question_type": "multiple_choice|true_false|fill_in_blank|descriptive",
      "options": ["A", "B", "C", "D"] or ["True", "False"] or [] for descriptive/fill-in-blank,
      "correct_answer": "The correct answer",
      "explanation": "Explanation of why this answer is correct",
      "topic": "Topic/category of the question",
      "difficulty": "{difficulty}"
    }}
  ]
}}

TEXT TO ANALYZE:
{text}

Generate exactly {num_questions} questions following the format above.
"""
    
    return base_prompt

def validate_enhanced_quiz_data(quiz_data, expected_questions, quiz_type):
    """Validate and clean quiz data for multiple question types"""
    if not isinstance(quiz_data, dict) or 'quiz' not in quiz_data:
        raise ValueError("Invalid quiz data format")
    
    questions = quiz_data['quiz']
    
    if not isinstance(questions, list):
        raise ValueError("Quiz questions must be a list")
    
    validated_questions = []
    
    for i, question in enumerate(questions):
        if len(validated_questions) >= expected_questions:
            break
            
        try:
            validated_question = validate_enhanced_question(question, quiz_type)
            validated_questions.append(validated_question)
        except ValueError as e:
            print(f"Skipping invalid question {i}: {str(e)}")
            continue
    
    if len(validated_questions) < expected_questions:
        raise ValueError(f"Only {len(validated_questions)} valid questions generated, expected {expected_questions}")
    
    return {"quiz": validated_questions}

def validate_enhanced_question(question, quiz_type):
    """Validate a single question with enhanced types"""
    required_fields = ['question_text', 'correct_answer']
    
    for field in required_fields:
        if field not in question:
            raise ValueError(f"Missing required field: {field}")
    
    # Set default values
    validated_question = {
        'question_text': question['question_text'],
        'question_type': question.get('question_type', 'multiple_choice'),
        'options': question.get('options', []),
        'correct_answer': question['correct_answer'],
        'explanation': question.get('explanation', ''),
        'topic': question.get('topic', 'General'),
        'difficulty': question.get('difficulty', 'medium')
    }
    
    # Validate based on question type
    if validated_question['question_type'] == 'multiple_choice':
        if not validated_question['options'] or len(validated_question['options']) < 2:
            raise ValueError("Multiple choice questions must have at least 2 options")
        
        if validated_question['correct_answer'] not in validated_question['options']:
            raise ValueError("Correct answer must be one of the options")
    
    elif validated_question['question_type'] == 'true_false':
        validated_question['options'] = ['True', 'False']
        if validated_question['correct_answer'] not in ['True', 'False']:
            raise ValueError("True/False questions must have 'True' or 'False' as correct answer")
    
    elif validated_question['question_type'] == 'fill_in_blank':
        validated_question['options'] = []
        if '___' not in validated_question['question_text']:
            # Add blank if not present
            validated_question['question_text'] += ' ___'
    
    elif validated_question['question_type'] == 'descriptive':
        validated_question['options'] = []
        # For descriptive, correct_answer can be a sample answer or key points
    
    return validated_question

# Keep existing helper functions
def clean_text(text):
    """Clean and preprocess extracted text"""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Remove special characters but keep basic punctuation
    text = re.sub(r'[^\w\s\.\,\?\!\;\:\-\(\)]', '', text)
    
    return text.strip()

def extract_json_from_response(content):
    """Extract JSON from OpenAI response"""
    try:
        # Try to parse the entire content as JSON
        return json.loads(content)
    except json.JSONDecodeError:
        # Look for JSON within code blocks
        json_pattern = r'``````'
        json_match = re.search(json_pattern, content, re.DOTALL)
        
        if json_match:
            json_str = json_match.group(1)
            return json.loads(json_str)
        
        # Look for JSON within the text
        json_pattern = r'\{.*\}'
        json_match = re.search(json_pattern, content, re.DOTALL)
        
        if json_match:
            json_str = json_match.group(0)
            return json.loads(json_str)
        
        raise ValueError("No valid JSON found in OpenAI response")
