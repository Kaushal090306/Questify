import json
import re
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from django.conf import settings
import requests
import google.generativeai as genai

def parse_document(file):
    ext = file.name.lower().split('.')[-1]
    if ext == 'pdf':
        return parse_pdf(file)
    elif ext == 'docx':
        return parse_docx(file)
    elif ext == 'pptx':
        return parse_pptx(file)
    else:
        raise ValueError("Unsupported file type.")

def parse_pdf(file):
    try:
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        text = "".join(page.get_text() for page in pdf)
        pdf.close()
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error parsing PDF: {str(e)}")

def parse_docx(file):
    try:
        doc = docx.Document(file)
        return "\n".join(p.text for p in doc.paragraphs).strip()
    except Exception as e:
        raise ValueError(f"Error parsing DOCX: {str(e)}")

def parse_pptx(file):
    try:
        ppt = Presentation(file)
        return "\n".join(
            shape.text
            for slide in ppt.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ).strip()
    except Exception as e:
        raise ValueError(f"Error parsing PPTX: {str(e)}")

def generate_quiz_with_gemini(text, num_questions=5, difficulty='medium', question_types=['multiple_choice'], topic='', document_text=''):
    # Handle both single quiz_type and list of question_types for backward compatibility
    if isinstance(question_types, str):
        quiz_type = question_types
    elif isinstance(question_types, list) and len(question_types) == 1:
        quiz_type = question_types[0]
    else:
        quiz_type = 'mixed'  # Use mixed when multiple types are selected
    
    max_retries = 3
    retry_count = 0
    attempted_minimal_prompt = False
    
    while retry_count < max_retries:
        try:
            print(f"DEBUG: Attempt {retry_count + 1} of {max_retries}")
            print(f"DEBUG: Checking Gemini API key...")
            if not settings.GOOGLE_GEMINI_API_KEY:
                raise ValueError("Google Gemini API key not configured in environment")
            print(f"DEBUG: API key found: {settings.GOOGLE_GEMINI_API_KEY[:10]}...")
            
            # Configure Gemini
            genai.configure(api_key=settings.GOOGLE_GEMINI_API_KEY)

            print(f"DEBUG: Cleaning text...")
            cleaned_text = clean_text(text)
            print(f"DEBUG: Original text length: {len(text)}, cleaned length: {len(cleaned_text)}")
            if len(cleaned_text) > 15000:
                cleaned_text = cleaned_text[:15000] + "..."
                print(f"DEBUG: Text truncated to 15000 characters")

            print(f"DEBUG: Creating prompt...")
            prompt = create_quiz_prompt(cleaned_text, num_questions, difficulty, quiz_type, question_types if isinstance(question_types, list) else None)
            print(f"DEBUG: Prompt length: {len(prompt)}")

            print(f"DEBUG: Initializing Gemini model...")
            model = genai.GenerativeModel('gemini-2.5-pro')
            print(f"DEBUG: Model initialized successfully")

            print(f"DEBUG: Generating content with Gemini...")
            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.5,
                    max_output_tokens=4096,
                )
            )
            print(f"DEBUG: Gemini response received")

            content = response.text
            print(f"DEBUG: Response content length: {len(content)}")
            print(f"DEBUG: Response preview: {content[:200]}...")

            print(f"DEBUG: Extracting JSON from response...")
            quiz_data = extract_json_from_response(content)
            print(f"DEBUG: JSON extracted successfully")

            print(f"DEBUG: Validating quiz data...")
            validated_quiz = validate_quiz_data(quiz_data, num_questions, quiz_type)
            print(f"DEBUG: Quiz data validated successfully")
            return validated_quiz
            
        except Exception as e:
            retry_count += 1
            print(f"DEBUG: Error in generate_quiz_with_openai (attempt {retry_count}): {str(e)}")
            
            # If no content or invalid JSON, try one simplified prompt attempt
            if ("No content" in str(e) or "No valid JSON" in str(e)) and not attempted_minimal_prompt:
                try:
                    print("DEBUG: Retrying once with minimal JSON-only prompt")
                    attempted_minimal_prompt = True
                    minimal_prompt = (
                        "Return ONLY valid JSON with key 'quiz' as an array of exactly "
                        f"{num_questions} items. Each item must include: 'question_text' (string), "
                        "'question_type' (one of multiple_choice, true_false, fill_in_blank, descriptive), "
                        "'options' (array, empty for non-multiple-choice), 'correct_answer' (string), "
                        "'explanation' (string), 'topic' (string), 'difficulty' (easy|medium|hard).\n\n"
                        "Base the questions on this content:\n" + clean_text(text)[:6000]
                    )
                    model = genai.GenerativeModel('gemini-2.5-pro')
                    response = model.generate_content(
                        minimal_prompt,
                        generation_config=genai.types.GenerationConfig(
                            temperature=0.3,
                            max_output_tokens=2048,
                        )
                    )
                    content = response.text
                    quiz_data = extract_json_from_response(content)
                    validated_quiz = validate_quiz_data(quiz_data, num_questions, quiz_type)
                    print("DEBUG: Minimal prompt succeeded")
                    return validated_quiz
                except Exception as me:
                    print(f"DEBUG: Minimal prompt also failed: {me}")

            # If it's a content blocking issue, don't retry
            if 'Content blocked' in str(e) or 'safety' in str(e).lower():
                print(f"DEBUG: Content blocked, not retrying")
                break
            
            # If we've exhausted retries, raise the error
            if retry_count >= max_retries:
                print(f"DEBUG: Max retries reached, raising error")
                break
            
            # Wait before retrying (respect API-provided retry_delay if present)
            import time
            wait_time = 2 ** retry_count
            msg = str(e)
            if '429' in msg:
                # Try to parse a suggested retry delay seconds: N
                import re
                m = re.search(r'retry_delay\s*\{\s*seconds:\s*(\d+)', msg)
                if m:
                    wait_time = max(wait_time, int(m.group(1)))
            print(f"DEBUG: Waiting {wait_time} seconds before retry...")
            time.sleep(wait_time)
    
    # If we get here, all retries failed - no fallback
    print(f"DEBUG: All attempts failed, no fallback available")
    raise Exception("Failed to generate quiz with Google Gemini after all retry attempts. Please try again later. [SERVICE_UNAVAILABLE]")

def create_quiz_prompt(text, num_questions, difficulty, quiz_type, question_types_list=None):
        diff_map = {
                'easy': 'basic understanding questions suitable for beginners',
                'medium': 'conceptual questions requiring moderate understanding',
                'hard': 'advanced questions testing deep comprehension'
        }
        type_map = {
                'multiple_choice': 'multiple choice questions with exactly 4 options (A, B, C, D)',
                'true_false': 'true/false questions',
                'fill_in_blank': 'fill-in-the-blank questions using ___ for missing words. The blank must be in the question text, and the correct answer must be provided separately.',
                'descriptive': 'short answer questions requiring brief explanations',
                'mixed': 'a mix of multiple choice, true/false, fill-in-the-blank (with ___ in the question text and correct answer provided), and descriptive questions. Do not repeat question types. At least one of each type if possible.'
        }

        # If specific question types are provided for mixed type, customize the prompt
        if quiz_type == 'mixed' and question_types_list:
            type_descriptions = []
            for q_type in question_types_list:
                if q_type in type_map:
                    type_descriptions.append(type_map[q_type])
            if type_descriptions:
                quiz_type_description = f"a mix of the following types: {', '.join(type_descriptions)}"
            else:
                quiz_type_description = type_map['mixed']
        else:
            quiz_type_description = type_map.get(quiz_type, type_map['mixed'])

        prompt = f"""
You are an educational AI assistant helping to create practice questions for students. Based on the provided educational material, create exactly {num_questions} {diff_map[difficulty]} questions of type {quiz_type_description}.

IMPORTANT:
- For fill-in-the-blank, use ___ in the question text for the blank, and provide the correct answer in the 'correct_answer' field.
- For mixed type, distribute question types evenly across the requested types.
- Return your response as a valid JSON object in this exact format:

{{
    "quiz": [
        {{
            "question_text": "What is the main concept discussed in this topic?",
            "question_type": "multiple_choice|true_false|fill_in_blank|descriptive",
            "options": ["Option A", "Option B", "Option C", "Option D"],  // empty for non-multiple-choice
            "correct_answer": "Option A or correct answer for blank",
            "explanation": "Brief explanation of why this answer is correct",
            "topic": "Main topic or concept",
            "difficulty": "{difficulty}"
        }}
    ]
}}

Educational material to base questions on:
{text}

Remember: Create questions that are educational, clear, and appropriate for academic learning.
"""

        return prompt

def clean_text(text):
    # Remove multiple whitespaces, keep punctuation
    text = re.sub(r'\s+', ' ', text)
    
    # Remove potentially problematic characters but keep educational content
    text = re.sub(r'[^\w\s\.\,\?\!\;\:\-\(\)\"\'\&\=\+\-\*\/\<\>\[\]]', '', text)
    
    # Remove very short lines that might contain noise
    lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 15]
    
    # Join lines and clean up
    cleaned = ' '.join(lines).strip()
    
    # Remove excessive punctuation
    cleaned = re.sub(r'[!]{2,}', '!', cleaned)
    cleaned = re.sub(r'[?]{2,}', '?', cleaned)
    cleaned = re.sub(r'[.]{2,}', '.', cleaned)
    
    # Ensure the text is not too long
    if len(cleaned) > 15000:
        # Try to find a good breaking point
        sentences = cleaned.split('.')
        truncated = ''
        for sentence in sentences:
            if len(truncated + sentence + '.') <= 15000:
                truncated += sentence + '.'
            else:
                break
        cleaned = truncated.strip()
    
    return cleaned

def extract_json_from_response(content):
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        match = re.search(r'\{[\s\S]*\}', content, re.DOTALL)
        if match:
            return json.loads(match.group(0))
        raise ValueError("No valid JSON found in Gemini response")

def validate_quiz_data(quiz_data, expected_questions, quiz_type):
    if not isinstance(quiz_data, dict) or 'quiz' not in quiz_data:
        raise ValueError("Invalid quiz data format: 'quiz' key missing")
    questions = quiz_data['quiz']
    if not isinstance(questions, list):
        raise ValueError("'quiz' key should contain a list")
    
    def normalize_question_type(raw_type: str) -> str:
        if not raw_type:
            return 'multiple_choice'
        t = str(raw_type).strip().lower()
        t = t.replace('-', '_').replace(' ', '_').replace('/', '_')
        # Common synonyms
        if t in {'mcq', 'mcqs', 'multiple_choice', 'multiple_choices'}:
            return 'multiple_choice'
        if t in {'true_false', 'true_or_false', 'boolean'}:
            return 'true_false'
        if t in {'fill_in_blank', 'fill_in_the_blank', 'fill_in_the_blanks', 'fill_in_blanks', 'fillintheblank', 'fill_the_blank'}:
            return 'fill_in_blank'
        if t in {'descriptive', 'short_answer', 'short_answers', 'open_ended', 'open_response', 'free_text'}:
            return 'descriptive'
        # For mixed type, we'll determine the specific type based on question content
        if t == 'mixed':
            return 'mixed'  # We'll handle this later
        # Default fallback
        return 'descriptive'
    def determine_question_type_from_content(question_data):
        """Determine the question type based on question content for mixed quizzes"""
        question_text = str(question_data.get('question_text', '')).lower()
        options = question_data.get('options', [])
        
        # Check for true/false questions
        if len(options) == 2 and set([str(o).lower().strip() for o in options]) == {'true', 'false'}:
            return 'true_false'
        
        # Check for fill in the blank questions
        if '___' in question_text or 'fill in' in question_text or 'blank' in question_text:
            return 'fill_in_blank'
        
        # Check for multiple choice questions (4 options)
        if len(options) >= 3:
            return 'multiple_choice'
        
        # Default to descriptive for open-ended questions
        return 'descriptive'
    
    validated_raw = []
    for q in questions[:expected_questions]:
        if not isinstance(q, dict):
            continue
        if 'question_text' not in q:
            continue
        normalized_type = normalize_question_type(q.get('question_type', quiz_type))
        
        # If the AI returned "mixed" as the question type, determine the actual type from content
        if normalized_type == 'mixed':
            normalized_type = determine_question_type_from_content(q)
        
        # Build base question
        question = {
            'question_text': str(q['question_text']).strip(),
            'question_type': normalized_type,
            'options': q.get('options', []) or [],
            'correct_answer': str(q.get('correct_answer', '')).strip(),
            'explanation': str(q.get('explanation', '') or '').strip(),
            'topic': str(q.get('topic', 'General') or 'General').strip(),
            'difficulty': q.get('difficulty', 'medium'),
        }
        validated_raw.append(question)

    # If mixed quiz type requested, distribute types among questions where needed
    if str(quiz_type).lower() == 'mixed':
        distribution = ['multiple_choice', 'true_false', 'fill_in_blank', 'descriptive']
        for idx, q in enumerate(validated_raw):
            if q['question_type'] not in distribution:
                q['question_type'] = distribution[idx % len(distribution)]
    
    # Final type-specific normalization of options
    validated = []
    for q in validated_raw:
        qtype = q['question_type']
        if qtype == 'multiple_choice':
            opts = [str(o) for o in (q.get('options') or [])]
            if len(opts) != 4:
                opts = ['Option A', 'Option B', 'Option C', 'Option D']
            q['options'] = opts
            if not q.get('correct_answer'):
                q['correct_answer'] = opts[0]
        elif qtype == 'true_false':
            q['options'] = ['True', 'False']
            if q.get('correct_answer') not in q['options']:
                q['correct_answer'] = 'True'
        else:
            q['options'] = []
            # For descriptive and fill_in_blank, allow empty correct_answer so UI can render
        validated.append(q)

    if not validated:
        raise ValueError("No valid questions could be extracted from response")
    return {"quiz": validated}
